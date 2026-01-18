import os
import json
import asyncio
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.types import InlineKeyboardMarkup, InlineKeyboardButton, FSInputFile
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocPt, RGBColor as DocRGB

# ============ DATABASE IMPORT ============
from database import get_db, init_db

# ============ KONFIGURATSIYA ============
BOT_TOKEN = "8434672153:AAE8IkxlJFBhRWtHwimy8l8UG272QtL4L30"
GEMINI_API_KEY = "AIzaSyB1Pn4PM5wZOlUJi4N4G29YQ49SOfCd5TM"
REQUIRED_CHANNEL = "@bkzsdfgahd"

# Gemini AI sozlash
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel('gemini-2.5-flash')

# Bot sozlash
bot = Bot(token=BOT_TOKEN)
storage = MemoryStorage()
dp = Dispatcher(storage=storage)

# Database instance
db = get_db()

# ============ HOLATLAR (STATES) ============
class BotStates(StatesGroup):
    lang_select = State()
    check_subscription = State()
    select_type = State()
    enter_topic = State()
    enter_pages = State()
    select_design = State()
    confirm = State()

# ============ TARJIMALAR ============
TEXTS = {
    'uz': {
        'welcome': '👋 Assalomu alaykum!\n\n📊 Taqdimot va 📝 Referat/Mustaqil ish tayyorlash botiga xush kelibsiz!\n\n💎 Kunlik limit: {remaining}/{total}\n\nTilni tanlang:',
        'select_lang': 'Tilni tanlang:',
        'subscription_required': '📢 Botdan foydalanish uchun kanalga obuna bo\'ling:\n\n{channel}\n\n✅ Obuna bo\'lgach "Tasdiqlash" tugmasini bosing',
        'check_btn': '✅ Obuna tekshirish',
        'not_subscribed': '❌ Siz hali obuna bo\'lmadingiz!\n\nIltimos, avval kanalga obuna bo\'ling: {channel}',
        'select_type': '📝 Qaysi turdagi hujjat kerak?\n\n💎 Bugungi limit: {remaining}/{total}',
        'presentation': '📊 Taqdimot (PPTX)',
        'report': '📝 Referat',
        'coursework': '📚 Mustaqil ish',
        'enter_topic': '✏️ Mavzuni kiriting:',
        'enter_pages': '📄 Nechta sahifa kerak? (3-50 oralig\'ida)',
        'invalid_pages': '❌ Noto\'g\'ri son! 3 dan 50 gacha son kiriting.',
        'select_design': '🎨 Dizayn shablonini tanlang:',
        'confirm_data': '📋 <b>Kiritilgan ma\'lumotlar:</b>\n\n'
                       '🎯 Tur: {doc_type}\n'
                       '📖 Mavzu: {topic}\n'
                       '📄 Sahifalar: {pages}\n'
                       '{design}'
                       '\n✅ Davom etamizmi?',
        'confirm_yes': '✅ Ha, davom etish',
        'confirm_no': '❌ Yo\'q, qaytadan',
        'generating': '⏳ Tayyorlanmoqda... Iltimos kuting...\n\n📊 Bu 30-60 soniya vaqt olishi mumkin.',
        'success': '✅ Tayyor! Marhamat:\n\n💎 Qolgan limit: {remaining}/{total}',
        'error': '❌ Xatolik yuz berdi. Iltimos qaytadan urinib ko\'ring.\n\nXatolik: {error}',
        'back_to_start': '🔙 Boshiga qaytish',
        'limit_reached': '⛔️ Kunlik limitingiz tugadi!\n\n'
                        '💎 Bugungi limit: {remaining}/{total}\n'
                        '🔄 Ertaga yangi limit beriladi\n\n'
                        '🎁 Ko\'proq limit olish uchun do\'stlaringizni taklif qiling!\n'
                        '🔗 Sizning referal havolangiz:\n'
                        '{ref_link}\n\n'
                        '👥 Har bir do\'stingiz uchun +1 doimiy limit!',
        'referral_success': '🎉 Tabriklaymiz!\n\n'
                           '{inviter} sizni taklif qildi va +1 limit oldi!\n\n'
                           '💎 Siz ham do\'stlaringizni taklif qilib limitingizni oshiring!',
        'referral_info': '👥 <b>Referal tizimi:</b>\n\n'
                        '💎 Sizning limitingiz: {limit}\n'
                        '📊 Taklif qilganlar: {count} ta\n'
                        '🔗 Referal havolangiz:\n'
                        '{ref_link}\n\n'
                        '🎁 Har bir do\'st uchun +1 doimiy limit!'
    },
    'ru': {
        'welcome': '👋 Здравствуйте!\n\n📊 Добро пожаловать в бот для создания презентаций и 📝 рефератов/курсовых работ!\n\n💎 Дневной лимит: {remaining}/{total}\n\nВыберите язык:',
        'select_lang': 'Выберите язык:',
        'subscription_required': '📢 Для использования бота подпишитесь на канал:\n\n{channel}\n\n✅ После подписки нажмите "Проверить"',
        'check_btn': '✅ Проверить подписку',
        'not_subscribed': '❌ Вы еще не подписались!\n\nПожалуйста, сначала подпишитесь на канал: {channel}',
        'select_type': '📝 Какой тип документа нужен?\n\n💎 Сегодняшний лимит: {remaining}/{total}',
        'presentation': '📊 Презентация (PPTX)',
        'report': '📝 Реферат',
        'coursework': '📚 Курсовая работа',
        'enter_topic': '✏️ Введите тему:',
        'enter_pages': '📄 Сколько страниц? (от 3 до 50)',
        'invalid_pages': '❌ Неверное число! Введите от 3 до 50.',
        'select_design': '🎨 Выберите дизайн шаблона:',
        'confirm_data': '📋 <b>Введенные данные:</b>\n\n'
                       '🎯 Тип: {doc_type}\n'
                       '📖 Тема: {topic}\n'
                       '📄 Страниц: {pages}\n'
                       '{design}'
                       '\n✅ Продолжить?',
        'confirm_yes': '✅ Да, продолжить',
        'confirm_no': '❌ Нет, заново',
        'generating': '⏳ Генерируется... Пожалуйста, подождите...\n\n📊 Это может занять 30-60 секунд.',
        'success': '✅ Готово! Держите:\n\n💎 Осталось: {remaining}/{total}',
        'error': '❌ Произошла ошибка. Попробуйте еще раз.\n\nОшибка: {error}',
        'back_to_start': '🔙 В начало',
        'limit_reached': '⛔️ Дневной лимит исчерпан!\n\n'
                        '💎 Сегодняшний лимит: {remaining}/{total}\n'
                        '🔄 Завтра получите новый лимит\n\n'
                        '🎁 Пригласите друзей для больше лимита!\n'
                        '🔗 Ваша реферальная ссылка:\n'
                        '{ref_link}\n\n'
                        '👥 За каждого друга +1 постоянный лимит!',
        'referral_success': '🎉 Поздравляем!\n\n'
                           '{inviter} пригласил вас и получил +1 лимит!\n\n'
                           '💎 Вы тоже можете пригласить друзей!',
        'referral_info': '👥 <b>Реферальная система:</b>\n\n'
                        '💎 Ваш лимит: {limit}\n'
                        '📊 Приглашено: {count} чел.\n'
                        '🔗 Реферальная ссылка:\n'
                        '{ref_link}\n\n'
                        '🎁 За каждого друга +1 постоянный лимит!'
    },
    'en': {
        'welcome': '👋 Hello!\n\n📊 Welcome to the presentation and 📝 report/coursework creation bot!\n\n💎 Daily limit: {remaining}/{total}\n\nSelect language:',
        'select_lang': 'Select language:',
        'subscription_required': '📢 To use the bot, subscribe to the channel:\n\n{channel}\n\n✅ After subscribing, click "Check"',
        'check_btn': '✅ Check subscription',
        'not_subscribed': '❌ You haven\'t subscribed yet!\n\nPlease subscribe to the channel first: {channel}',
        'select_type': '📝 What type of document do you need?\n\n💎 Today\'s limit: {remaining}/{total}',
        'presentation': '📊 Presentation (PPTX)',
        'report': '📝 Report',
        'coursework': '📚 Coursework',
        'enter_topic': '✏️ Enter topic:',
        'enter_pages': '📄 How many pages? (3 to 50)',
        'invalid_pages': '❌ Invalid number! Enter from 3 to 50.',
        'select_design': '🎨 Select design template:',
        'confirm_data': '📋 <b>Entered data:</b>\n\n'
                       '🎯 Type: {doc_type}\n'
                       '📖 Topic: {topic}\n'
                       '📄 Pages: {pages}\n'
                       '{design}'
                       '\n✅ Continue?',
        'confirm_yes': '✅ Yes, continue',
        'confirm_no': '❌ No, restart',
        'generating': '⏳ Generating... Please wait...\n\n📊 This may take 30-60 seconds.',
        'success': '✅ Done! Here you go:\n\n💎 Remaining: {remaining}/{total}',
        'error': '❌ An error occurred. Please try again.\n\nError: {error}',
        'back_to_start': '🔙 Back to start',
        'limit_reached': '⛔️ Daily limit reached!\n\n'
                        '💎 Today\'s limit: {remaining}/{total}\n'
                        '🔄 New limit tomorrow\n\n'
                        '🎁 Invite friends for more limit!\n'
                        '🔗 Your referral link:\n'
                        '{ref_link}\n\n'
                        '👥 +1 permanent limit for each friend!',
        'referral_success': '🎉 Congratulations!\n\n'
                           '{inviter} invited you and got +1 limit!\n\n'
                           '💎 You can also invite friends!',
        'referral_info': '👥 <b>Referral system:</b>\n\n'
                        '💎 Your limit: {limit}\n'
                        '📊 Invited: {count} people\n'
                        '🔗 Referral link:\n'
                        '{ref_link}\n\n'
                        '🎁 +1 permanent limit for each friend!'
    }
}

# ============ DIZAYN SHABLONLARI ============
DESIGNS = {
    '1': {'name': 'Klassik Ko\'k', 'bg': (31, 78, 121), 'title': (255, 255, 255), 'text': (0, 0, 0)},
    '2': {'name': 'Professional', 'bg': (68, 114, 196), 'title': (255, 255, 255), 'text': (0, 0, 0)},
    '3': {'name': 'Zamonaviy', 'bg': (91, 155, 213), 'title': (255, 255, 255), 'text': (0, 0, 0)},
    '4': {'name': 'Qizil energiya', 'bg': (192, 0, 0), 'title': (255, 255, 255), 'text': (0, 0, 0)},
    '5': {'name': 'Yashil tabiat', 'bg': (0, 176, 80), 'title': (255, 255, 255), 'text': (0, 0, 0)}
}

# ============ YORDAMCHI FUNKSIYALAR ============
async def check_subscription(user_id: int) -> bool:
    """Foydalanuvchi obuna holatini tekshirish"""
    try:
        member = await bot.get_chat_member(chat_id=REQUIRED_CHANNEL, user_id=user_id)
        return member.status in ['member', 'administrator', 'creator']
    except:
        return False

def get_text(lang: str, key: str) -> str:
    """Tanlangan tildagi matnni olish"""
    return TEXTS.get(lang, TEXTS['uz']).get(key, '')

# ============ GEMINI AI ORQALI KONTENT OLISH ============
async def generate_content_with_gemini(topic: str, pages: int, doc_type: str, lang: str) -> dict:
    """Gemini AI yordamida kontent generatsiya qilish"""
    
    lang_map = {'uz': 'uzbek', 'ru': 'russian', 'en': 'english'}
    lang_full = lang_map.get(lang, 'uzbek')
    
    if doc_type == 'presentation':
        prompt = f"""Create a detailed presentation content in {lang_full} language about "{topic}".
        
Generate EXACTLY {pages} slides with the following structure:

Return ONLY valid JSON (no markdown, no extra text) in this exact format:
{{
  "title": "Main presentation title",
  "slides": [
    {{
      "title": "Slide 1 title",
      "content": [
        "First point about the topic",
        "Second point with details",
        "Third important point"
      ]
    }},
    ... (continue for all {pages} slides)
  ]
}}

Requirements:
- Each slide must have 3-5 bullet points
- Content must be informative and well-structured
- Use {lang_full} language throughout
- Make it educational and professional
- RETURN ONLY JSON, NO OTHER TEXT"""

    else:  # report or coursework
        prompt = f"""Create a detailed {'report' if doc_type == 'report' else 'coursework'} in {lang_full} language about "{topic}".

Generate content for approximately {pages} pages with the following structure:

Return ONLY valid JSON (no markdown, no extra text) in this exact format:
{{
  "title": "Document title",
  "introduction": "Detailed introduction (2-3 paragraphs)",
  "sections": [
    {{
      "title": "Section 1 title",
      "content": "Detailed content for this section (3-4 paragraphs)"
    }},
    {{
      "title": "Section 2 title",
      "content": "Detailed content for this section (3-4 paragraphs)"
    }},
    ... (continue with more sections)
  ],
  "conclusion": "Detailed conclusion (2-3 paragraphs)"
}}

Requirements:
- Create enough sections to fill {pages} pages
- Each section should have detailed, informative content
- Use {lang_full} language throughout
- Make it academic and well-researched
- RETURN ONLY JSON, NO OTHER TEXT"""

    try:
        response = model.generate_content(prompt)
        result_text = response.text.strip()
        
        # Agar markdown formatida kelsa, tozalash
        if result_text.startswith('```'):
            result_text = result_text.split('```')[1]
            if result_text.startswith('json'):
                result_text = result_text[4:]
            result_text = result_text.strip()
        
        data = json.loads(result_text)
        return data
    except json.JSONDecodeError as e:
        print(f"JSON parse error: {e}")
        print(f"Response text: {result_text[:500]}")
        raise Exception("AI javobini tahlil qilishda xatolik")
    except Exception as e:
        print(f"Gemini error: {e}")
        raise

# ============ PPTX YARATISH ============
def create_presentation(data: dict, design_id: str, output_path: str):
    """PPTX taqdimot yaratish"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    design = DESIGNS[design_id]
    bg_color = design['bg']
    title_color = design['title']
    text_color = design['text']
    
    # Birinchi slayd - sarlavha
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)
    
    # Sarlavha
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = data['title']
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_color)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Qolgan slaydlar
    for slide_data in data['slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Sarlavha
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*bg_color)
        
        # Kontent
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for point in slide_data['content']:
            p = text_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*text_color)
            p.space_before = Pt(12)
            p.level = 0
    
    prs.save(output_path)

# ============ WORD YARATISH ============
def create_document(data: dict, output_path: str):
    """Word hujjat yaratish"""
    doc = Document()
    
    # Sarlavha
    title = doc.add_heading(data['title'], 0)
    title.alignment = 1  # Center
    
    # Kirish
    doc.add_heading('Kirish', 1)
    doc.add_paragraph(data['introduction'])
    
    # Bo'limlar
    for section in data['sections']:
        doc.add_heading(section['title'], 1)
        doc.add_paragraph(section['content'])
    
    # Xulosa
    doc.add_heading('Xulosa', 1)
    doc.add_paragraph(data['conclusion'])
    
    doc.save(output_path)

# ============ HANDLERLAR ============
@dp.message(Command("start"))
async def cmd_start(message: types.Message, state: FSMContext):
    """Start buyrug'i - referal tizimi bilan"""
    await state.clear()
    
    user_id = message.from_user.id
    username = message.from_user.username
    first_name = message.from_user.first_name
    
    # Foydalanuvchini database'ga qo'shish yoki olish
    db.create_user(user_id, username, first_name)
    
    # Referal tekshirish
    args = message.text.split()
    if len(args) > 1:
        try:
            referrer_id = int(args[1])
            if referrer_id != user_id:
                # Yangi foydalanuvchi referral orqali kelgan
                success = db.add_referral(referrer_id, user_id)
                
                if success:
                    # Taklif qilgan foydalanuvchiga xabar
                    try:
                        referrer = db.get_user(referrer_id)
                        ref_lang = referrer['language'] if referrer else 'uz'
                        
                        await bot.send_message(
                            referrer_id,
                            f"🎉 Yangi foydalanuvchi sizning havolangiz orqali botga qo'shildi!\n\n"
                            f"💎 Sizga +1 doimiy limit berildi!\n"
                            f"📊 Jami limitingiz: {referrer['daily_limit'] + 1}"
                        )
                    except:
                        pass
        except:
            pass
    
    # Limit ma'lumotlarini olish
    remaining, total = db.get_daily_limit(user_id)
    
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🇺🇿 O'zbek", callback_data="lang_uz")],
        [InlineKeyboardButton(text="🇷🇺 Русский", callback_data="lang_ru")],
        [InlineKeyboardButton(text="🇬🇧 English", callback_data="lang_en")]
    ])
    
    welcome_text = (
        "👋 Assalomu alaykum! | Здравствуйте! | Hello!\n\n"
        "📊 Taqdimot va 📝 Referat/Mustaqil ish tayyorlash boti\n\n"
        f"💎 Kunlik limit: {remaining}/{total}\n\n"
        "Tilni tanlang | Выберите язык | Select language:"
    )
    
    await message.answer(welcome_text, reply_markup=keyboard)
    await state.set_state(BotStates.lang_select)

@dp.message(Command("referral"))
async def cmd_referral(message: types.Message, state: FSMContext):
    """Referal ma'lumotlarini ko'rish"""
    user_id = message.from_user.id
    
    # Get user from database
    user = db.get_user(user_id)
    if not user:
        await message.answer("⚠️ Iltimos, avval /start ni bosing")
        return
    
    lang = user['language']
    
    # Get referral count from database
    ref_count = db.get_referral_count(user_id)
    bot_username = (await bot.me()).username
    ref_link = f"https://t.me/{bot_username}?start={user_id}"
    
    text = get_text(lang, 'referral_info').format(
        limit=user['daily_limit'],
        count=ref_count,
        ref_link=ref_link
    )
    
    await message.answer(text, parse_mode='HTML')

@dp.callback_query(F.data.startswith("lang_"))
async def process_language(callback: types.CallbackQuery, state: FSMContext):
    """Til tanlash"""
    lang = callback.data.split("_")[1]
    await state.update_data(lang=lang)
    
    # Update language in database
    user_id = callback.from_user.id
    db.users.update_language(user_id, lang)
    
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=get_text(lang, 'check_btn'), callback_data="check_sub")]
    ])
    
    text = get_text(lang, 'subscription_required').format(channel=REQUIRED_CHANNEL)
    await callback.message.edit_text(text, reply_markup=keyboard)
    await state.set_state(BotStates.check_subscription)
    await callback.answer()

@dp.callback_query(F.data == "check_sub")
async def check_sub(callback: types.CallbackQuery, state: FSMContext):
    """Obuna tekshirish"""
    user_id = callback.from_user.id
    user = db.get_user(user_id)
    lang = user['language'] if user else 'uz'
    
    is_subscribed = await check_subscription(user_id)
    
    if is_subscribed:
        # Limit tekshirish
        remaining, total = db.get_daily_limit(user_id)
        
        keyboard = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text=get_text(lang, 'presentation'), callback_data="type_presentation")],
            [InlineKeyboardButton(text=get_text(lang, 'report'), callback_data="type_report")],
            [InlineKeyboardButton(text=get_text(lang, 'coursework'), callback_data="type_coursework")]
        ])
        
        text = get_text(lang, 'select_type').format(remaining=remaining, total=total)
        
        await callback.message.edit_text(text, reply_markup=keyboard)
        await state.set_state(BotStates.select_type)
    else:
        await callback.answer(
            get_text(lang, 'not_subscribed').format(channel=REQUIRED_CHANNEL),
            show_alert=True
        )
    
    await callback.answer()

@dp.callback_query(F.data.startswith("type_"))
async def select_type(callback: types.CallbackQuery, state: FSMContext):
    """Hujjat turini tanlash"""
    user_id = callback.from_user.id
    user = db.get_user(user_id)
    lang = user['language'] if user else 'uz'
    
    # Limit tekshirish
    if not db.can_generate(user_id):
        bot_username = (await bot.me()).username
        ref_link = f"https://t.me/{bot_username}?start={user_id}"
        remaining, total = db.get_daily_limit(user_id)
        
        text = get_text(lang, 'limit_reached').format(
            remaining=remaining,
            total=total,
            ref_link=ref_link
        )
        
        await callback.message.edit_text(text, parse_mode='HTML')
        await callback.answer()
        return
    
    doc_type = callback.data.split("_")[1]
    await state.update_data(doc_type=doc_type)
    
    await callback.message.edit_text(get_text(lang, 'enter_topic'))
    await state.set_state(BotStates.enter_topic)
    await callback.answer()

@dp.message(BotStates.enter_topic)
async def enter_topic(message: types.Message, state: FSMContext):
    """Mavzu kiritish"""
    await state.update_data(topic=message.text)
    
    user = db.get_user(message.from_user.id)
    lang = user['language'] if user else 'uz'
    
    await message.answer(get_text(lang, 'enter_pages'))
    await state.set_state(BotStates.enter_pages)

@dp.message(BotStates.enter_pages)
async def enter_pages(message: types.Message, state: FSMContext):
    """Sahifalar sonini kiritish"""
    user = db.get_user(message.from_user.id)
    lang = user['language'] if user else 'uz'
    
    try:
        pages = int(message.text)
        if pages < 3 or pages > 50:
            await message.answer(get_text(lang, 'invalid_pages'))
            return
        
        await state.update_data(pages=pages)
        data = await state.get_data()
        doc_type = data.get('doc_type')
        
        if doc_type == 'presentation':
            # Dizayn tanlash
            designs_text = "🎨 <b>Mavjud dizaynlar:</b>\n\n"
            keyboard_buttons = []
            
            for design_id, design_info in DESIGNS.items():
                designs_text += f"{design_id}. {design_info['name']}\n"
                keyboard_buttons.append([
                    InlineKeyboardButton(text=f"{design_id}. {design_info['name']}", callback_data=f"design_{design_id}")
                ])
            
            keyboard = InlineKeyboardMarkup(inline_keyboard=keyboard_buttons)
            await message.answer(designs_text, reply_markup=keyboard, parse_mode='HTML')
            await state.set_state(BotStates.select_design)
        else:
            # To'g'ridan-to'g'ri tasdiqga o'tish
            await show_confirmation(message, state)
            
    except ValueError:
        await message.answer(get_text(lang, 'invalid_pages'))

@dp.callback_query(F.data.startswith("design_"))
async def select_design(callback: types.CallbackQuery, state: FSMContext):
    """Dizayn tanlash"""
    design_id = callback.data.split("_")[1]
    await state.update_data(design=design_id)
    await show_confirmation(callback.message, state)
    await callback.answer()

async def show_confirmation(message: types.Message, state: FSMContext):
    """Ma'lumotlarni tasdiqlash"""
    data = await state.get_data()
    user = db.get_user(message.from_user.id if hasattr(message, 'from_user') else message.chat.id)
    lang = user['language'] if user else 'uz'
    
    doc_type = data.get('doc_type')
    topic = data.get('topic')
    pages = data.get('pages')
    design = data.get('design')
    
    type_names = {
        'uz': {'presentation': 'Taqdimot', 'report': 'Referat', 'coursework': 'Mustaqil ish'},
        'ru': {'presentation': 'Презентация', 'report': 'Реферат', 'coursework': 'Курсовая работа'},
        'en': {'presentation': 'Presentation', 'report': 'Report', 'coursework': 'Coursework'}
    }
    
    doc_type_name = type_names[lang][doc_type]
    design_text = f"🎨 Dizayn: {DESIGNS[design]['name']}\n" if design else ""
    
    text = get_text(lang, 'confirm_data').format(
        doc_type=doc_type_name,
        topic=topic,
        pages=pages,
        design=design_text
    )
    
    keyboard = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text=get_text(lang, 'confirm_yes'), callback_data="confirm_yes")],
        [InlineKeyboardButton(text=get_text(lang, 'confirm_no'), callback_data="confirm_no")]
    ])
    
    await message.answer(text, reply_markup=keyboard, parse_mode='HTML')
    await state.set_state(BotStates.confirm)

@dp.callback_query(F.data == "confirm_yes")
async def confirm_yes(callback: types.CallbackQuery, state: FSMContext):
    """Tasdiqlash - hujjat yaratish"""
    data = await state.get_data()
    user_id = callback.from_user.id
    user = db.get_user(user_id)
    lang = user['language'] if user else 'uz'
    
    doc_type = data.get('doc_type')
    topic = data.get('topic')
    pages = data.get('pages')
    design = data.get('design')
    
    await callback.message.edit_text(get_text(lang, 'generating'))
    
    # Create generation record in database
    generation_id = db.generations.create_generation(user_id, doc_type, topic, pages, design)
    
    try:
        # AI dan kontent olish
        content = await generate_content_with_gemini(topic, pages, doc_type, lang)
        
        # Fayl yaratish
        if doc_type == 'presentation':
            filename = f"presentation_{user_id}_{generation_id}.pptx"
            create_presentation(content, design, filename)
        else:
            filename = f"document_{user_id}_{generation_id}.docx"
            create_document(content, filename)
        
        # Update generation status to completed
        db.generations.update_status(generation_id, 'completed', filename)
        
        # Use generation (decrease limit)
        db.use_generation(user_id)
        
        # Get updated limits
        remaining, total = db.get_daily_limit(user_id)
        
        # Faylni yuborish
        file = FSInputFile(filename)
        await callback.message.answer_document(
            document=file,
            caption=get_text(lang, 'success').format(remaining=remaining, total=total)
        )
        
        # Faylni o'chirish
        if os.path.exists(filename):
            os.remove(filename)
        
        # Boshiga qaytish tugmasi
        keyboard = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text=get_text(lang, 'back_to_start'), callback_data="back_start")]
        ])
        await callback.message.answer("👍", reply_markup=keyboard)
        
        await state.clear()
        
    except Exception as e:
        # Update generation status to failed
        db.generations.update_status(generation_id, 'failed', error_message=str(e))
        
        error_text = get_text(lang, 'error').format(error=str(e))
        await callback.message.answer(error_text)
        print(f"Xatolik yuz berdi: {e}")
        import traceback
        traceback.print_exc()
    
    await callback.answer()

@dp.callback_query(F.data == "confirm_no")
async def confirm_no(callback: types.CallbackQuery, state: FSMContext):
    """Bekor qilish - qaytadan boshlash"""
    await state.clear()
    await cmd_start(callback.message, state)
    await callback.answer()

@dp.callback_query(F.data == "back_start")
async def back_to_start(callback: types.CallbackQuery, state: FSMContext):
    """Boshiga qaytish"""
    await state.clear()
    await cmd_start(callback.message, state)
    await callback.answer()

# ============ BOTNI ISHGA TUSHIRISH ============
async def main():
    """Botni ishga tushirish funksiyasi"""
    print("=" * 50)
    print("🤖 Bot ishga tushmoqda...")
    print("=" * 50)
    print(f"📊 Bot nomi: Presentation & Report Generator")
    print(f"🔑 Token: {BOT_TOKEN[:10]}...")
    print(f"📢 Majburiy kanal: {REQUIRED_CHANNEL}")
    print(f"🤖 AI Model: Gemini Pro")
    print(f"💾 Database: SQLite (bot_data.db)")
    print("=" * 50)
    
    # Initialize database
    try:
        init_db()
        print("✅ Database initialized successfully!")
    except Exception as e:
        print(f"⚠️ Database initialization warning: {e}")
    
    print("✅ Bot muvaffaqiyatli ishga tushdi!")
    print("💬 Xabarlarni kutmoqda...\n")
    
    try:
        await dp.start_polling(bot, skip_updates=True)
    except Exception as e:
        print(f"❌ Bot ishga tushirishda xatolik: {e}")

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        print("\n\n🛑 Bot to'xtatildi!")
    except Exception as e:
        print(f"\n❌ Kritik xatolik: {e}")
        import traceback
        traceback.print_exc()